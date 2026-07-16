#!/usr/bin/env python3
"""Create editable Fig. 1: problem-to-method workflow."""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "intro_problem_method_workflow"
prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(2.15)
s = prs.slides.add_slide(prs.slide_layouts[6])
NAVY, BLUE, ORANGE, GREEN, GRAY = RGBColor(24,73,122), RGBColor(220,238,250), RGBColor(252,235,218), RGBColor(222,242,232), RGBColor(70,70,70)
def box(x,y,w,h,t,c):
 sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.color.rgb=NAVY; sh.line.width=Pt(.8)
 tf=sh.text_frame; tf.clear(); tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.name="Arial"; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=GRAY
def arrow(x1,y1,x2,y2):
 c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),Inches(y1),Inches(x2),Inches(y2)); c.line.color.rgb=NAVY; c.line.width=Pt(1.4); c.line.end_arrowhead=True
box(.35,.45,2.65,1.2,"Problem\nSingle-carrier readout limits\ncarrier-domain parallelism",ORANGE)
box(3.55,.45,2.85,1.2,"Method\nFour carrier-local PMOS DACs\n+ shared RRAM MVM",BLUE)
box(6.95,.45,2.85,1.2,"Physical recovery\nTIA + four I/Q receiver pairs\n+ sampling / decision",BLUE)
box(10.35,.45,2.65,1.2,"Evidence boundary\nTransient peripheral check\n+ mapped-network evaluation",GREEN)
arrow(3.0,1.05,3.55,1.05); arrow(6.4,1.05,6.95,1.05); arrow(9.8,1.05,10.35,1.05)
for x,t in [(.35,"Carrier-concurrent input encoding"),(3.55,"Physical modulation and coherent recovery"),(7.15,"Separate circuit and array-scale claims")]:
 q=s.shapes.add_textbox(Inches(x),Inches(1.78), Inches(3),Inches(.2)); p=q.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=t; r.font.name="Arial"; r.font.size=Pt(7.5); r.font.color.rgb=GRAY
prs.save(OUT.with_suffix('.pptx'))
