#!/usr/bin/env python3
"""Create an editable physical four-carrier architecture figure (PPTX)."""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "physical_fourcarrier_architecture"
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

NAVY, BLUE, LIGHTBLUE = RGBColor(24, 73, 122), RGBColor(32, 110, 170), RGBColor(222, 239, 250)
ORANGE, LIGHTORANGE = RGBColor(194, 93, 0), RGBColor(252, 235, 218)
GREEN, LIGHTGREEN = RGBColor(42, 127, 98), RGBColor(222, 242, 232)
GRAY, LIGHTGRAY, WHITE = RGBColor(80, 80, 80), RGBColor(240, 242, 244), RGBColor(255, 255, 255)

def box(x, y, w, h, text, fill, line=NAVY, size=10, bold=False, radius=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line; shape.line.width = Pt(0.8)
    frame = shape.text_frame; frame.clear(); frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text; run.font.name = "Arial"; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = GRAY
    return shape

def label(x, y, w, h, text, size=10, color=GRAY, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame; frame.clear(); frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return shape

def arrow(x1, y1, x2, y2, color=NAVY, width=1.2):
    con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    con.line.color.rgb = color; con.line.width = Pt(width); con.line.end_arrowhead = True
    return con

label(0.35, 0.1, 12.6, 0.36, "Physical Four-Carrier Modulation and Readout Architecture", 17, NAVY, True)
label(0.35, 0.46, 12.6, 0.26, "45 nm compact-model transient implementation; 128×128 mapped crossbar is evaluated separately", 8.5, GRAY)

box(0.35, 1.15, 1.25, 0.8, "2-bit input\n00 / 01 / 10 / 11", LIGHTGRAY, GRAY, 10, True, True)
box(1.9, 0.95, 1.55, 1.2, "Four PMOS\ncurrent-cell DACs\n250/500/750 MHz/1 GHz", LIGHTORANGE, ORANGE, 9, True, True)
box(1.9, 2.35, 1.55, 0.72, "Ring LO +\n90° phase network", LIGHTBLUE, BLUE, 9, True, True)
box(3.78, 1.15, 1.35, 0.8, "Wordline\nsumming driver", LIGHTORANGE, ORANGE, 10, True, True)
box(5.52, 0.88, 1.35, 1.35, "RRAM crossbar\n128 × 128\nconductance matrix", LIGHTGREEN, GREEN, 10, True, True)
box(7.25, 1.15, 1.15, 0.8, "Selected BL\nTIA", LIGHTORANGE, ORANGE, 10, True, True)

arrow(1.60, 1.55, 1.90, 1.55); arrow(3.45, 1.55, 3.78, 1.55); arrow(5.13, 1.55, 5.52, 1.55); arrow(6.87, 1.55, 7.25, 1.55)
arrow(2.68, 2.35, 2.68, 2.15, BLUE); label(1.78, 3.1, 1.8, 0.28, "shared LO/reference distribution", 8, BLUE)

label(8.7, 0.78, 3.9, 0.26, "Four carrier-selective I/Q receiver pairs (8 mixers)", 10, NAVY, True)
freqs = ["250 MHz", "500 MHz", "750 MHz", "1 GHz"]
for i, freq in enumerate(freqs):
    y = 1.1 + i * 1.08
    box(8.72, y, 0.82, 0.72, freq, LIGHTBLUE, BLUE, 8.3, True, True)
    box(9.78, y, 1.16, 0.72, "I/Q Gilbert\nmixers", LIGHTBLUE, BLUE, 8.3, False, True)
    box(11.18, y, 1.16, 0.72, "LP integration\n+ S/H", LIGHTGRAY, GRAY, 8.3, False, True)
    box(12.56, y, 0.48, 0.72, "1b\nD", LIGHTORANGE, ORANGE, 8.3, True, True)
    arrow(8.40, 1.55, 8.72, y + 0.36, ORANGE, 0.8)
    arrow(9.54, y + 0.36, 9.78, y + 0.36); arrow(10.94, y + 0.36, 11.18, y + 0.36); arrow(12.34, y + 0.36, 12.56, y + 0.36)

box(8.72, 5.75, 1.55, 0.74, "Symbol\nde-mapping", LIGHTGRAY, GRAY, 9, True, True)
box(10.58, 5.75, 1.55, 0.74, "Digital accumulation\n+ activation", LIGHTGRAY, GRAY, 9, True, True)
box(12.38, 5.75, 0.66, 0.74, "Layer\ncontrol", LIGHTGRAY, GRAY, 8, True, True)
arrow(12.8, 5.45, 9.48, 5.75, GRAY, 0.8); arrow(10.27, 6.12, 10.58, 6.12, GRAY); arrow(12.13, 6.12, 12.38, 6.12, GRAY)
label(0.45, 6.86, 12.4, 0.28, "Solid blocks: transistor-level peripheral implementation. Crossbar/network mapping: separate system-level evaluation with the same carrier organization.", 8.3, GRAY)
prs.save(OUT.with_suffix(".pptx"))
